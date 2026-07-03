"""Generate the LearnStack investor pitch deck (PITCH_DECK.pptx).

Run from repo root:  backend/venv/Scripts/python scripts/make_pitch_deck.py
Edit the CONTENT section and re-run to regenerate. Numbers marked EDIT
(fundraise amount, use of funds) are placeholders for the founder to set.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

GREEN = RGBColor(0x1A, 0x7F, 0x5A)
NAVY = RGBColor(0x12, 0x32, 0x5C)
DARK = RGBColor(0x1F, 0x24, 0x28)
GREY = RGBColor(0x5F, 0x6B, 0x73)
LIGHT = RGBColor(0xF2, 0xF5, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size, color=DARK, bold=False, first=False, align=PP_ALIGN.LEFT,
         space_after=6, font="Segoe UI"):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    f = r.font
    f.size, f.bold, f.name = Pt(size), bold, font
    f.color.rgb = color
    return p


def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def header(s, kicker, title):
    rect(s, 0, 0, W, Inches(0.18), GREEN)
    tf = box(s, Inches(0.7), Inches(0.45), Inches(12), Inches(1.4))
    para(tf, kicker.upper(), 13, GREEN, bold=True, first=True)
    para(tf, title, 33, NAVY, bold=True)


def bullets(s, items, x=Inches(0.7), y=Inches(1.95), w=Inches(11.9), size=17):
    tf = box(s, x, y, w, H - y - Inches(0.5))
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            lead, rest = item
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(11)
            r1 = p.add_run(); r1.text = f"{lead}  "
            r1.font.size, r1.font.bold, r1.font.name = Pt(size), True, "Segoe UI"
            r1.font.color.rgb = NAVY
            r2 = p.add_run(); r2.text = rest
            r2.font.size, r2.font.name = Pt(size), "Segoe UI"
            r2.font.color.rgb = DARK
        else:
            para(tf, item, size, DARK, first=(i == 0), space_after=11)


def stat_row(s, stats, y=Inches(4.9), height=Inches(1.7)):
    n = len(stats)
    gap = Inches(0.25)
    total_gap = Emu(int(gap) * (n - 1))
    card_w = Emu(int(W - Inches(1.4) - total_gap) // n)
    x = Inches(0.7)
    for value, label in stats:
        card = rect(s, x, y, card_w, height, LIGHT)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, value, 30, GREEN, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=2)
        para(tf, label, 12.5, GREY, align=PP_ALIGN.CENTER)
        x = Emu(int(x) + int(card_w) + int(gap))


def table_slide(s, headers_row, rows, y=Inches(1.95), col_widths=None):
    n_cols = len(headers_row)
    x, w = Inches(0.7), Inches(11.9)
    tbl = s.shapes.add_table(len(rows) + 1, n_cols, x, y, w,
                             Inches(0.42) * (len(rows) + 1)).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    def set_cell(cell, text, size, color, bold, bg):
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        pr = cell.text_frame.paragraphs[0]
        run = pr.runs[0] if pr.runs else pr.add_run()
        run.font.size, run.font.bold, run.font.name = Pt(size), bold, "Segoe UI"
        run.font.color.rgb = color

    for c, text in enumerate(headers_row):
        set_cell(tbl.cell(0, c), text, 13, WHITE, True, NAVY)
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            set_cell(tbl.cell(r, c), text, 12.5, DARK, c == 0,
                     WHITE if r % 2 else LIGHT)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ============================== CONTENT ==============================

# 1 — Title
s = slide(NAVY)
rect(s, 0, Inches(6.9), W, Inches(0.6), GREEN)
tf = box(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(3.5))
para(tf, "LearnStack", 60, WHITE, bold=True, first=True, space_after=10)
para(tf, "White-label curriculum delivery, as a service", 26, RGBColor(0xBF, 0xD9, 0xCE), space_after=22)
para(tf, "Shopify made everyone a store.  LearnStack makes everyone an academy.",
     18, WHITE, bold=True)
tf2 = box(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.6))
para(tf2, "Investor briefing · July 2026 · Lagos, Nigeria", 13, RGBColor(0x9F, 0xB4, 0xC6), first=True)
notes(s, "Add founder name + contact here.")

# 2 — Problem
s = slide()
header(s, "The problem", "Teaching at scale means choosing between three bad options")
bullets(s, [
    ("Build your own platform.", "6+ months and an engineering team most training businesses, NGOs, and agencies don't have."),
    ("Rent Teachable / Thinkific.", "Fast — but the academy wears someone else's brand, pricing punishes growth, and content is locked in."),
    ("Self-host Moodle.", "Free software, expensive reality: servers, plugins, and an interface learners tolerate rather than enjoy."),
    ("And every one of them ignores the biggest audience:", "learners on WhatsApp, basic phones, and low bandwidth — ~60% of Nigeria can't be reached by smartphone-only platforms."),
])
notes(s, "Anchor on the customer segments: academies, government programs, corporate L&D, edtech founders.")

# 3 — Solution
s = slide()
header(s, "The solution", "A branded academy in a day — delivering any curriculum, on any channel")
bullets(s, [
    ("White-label multi-tenancy.", "Each customer gets their name, logo, colors, and domain on one shared platform — unlimited brands, one codebase."),
    ("Author once, in blocks.", "Lessons combine video, rich text, audio, images, files, embeds, and graded quizzes."),
    ("Deliver everywhere.", "Full-fidelity web · WhatsApp with real video/audio messages · SMS · USSD for feature phones."),
    ("Prove it.", "Progress tracking and certificates any employer or ministry can verify publicly."),
], size=17)
stat_row(s, [("1 day", "brand to launch"), ("7 formats", "per lesson"), ("4 channels", "web · WhatsApp · SMS · USSD"), ("100%", "curriculum portability")])

# 4 — Product / how it works
s = slide()
header(s, "Product", "One lesson, rendered to each channel's strengths")
table_slide(s, ["Channel", "What the learner gets"], [
    ["Web", "Streaming video, audio, quizzes, certificates — the full experience"],
    ["WhatsApp", "Conversational lessons with real video, audio, and image messages; quiz by reply"],
    ["SMS", "Text lessons chunked to segment limits; media becomes a captioned link"],
    ["USSD", "Text-only sessions for feature phones — no data plan, no smartphone required"],
], col_widths=[Inches(2.4), Inches(9.5)])
tf = box(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(2.3))
para(tf, "Under the hood", 14, GREEN, bold=True, first=True, space_after=4)
para(tf, "One conversation engine drives every messaging channel: menu → enroll → lesson → quiz → certificate. "
         "Tenants connect their own WhatsApp/SMS numbers. A built-in simulator tests every flow before "
         "any provider account exists. Whole curricula import and export as portable JSON.", 15, DARK)

# 5 — Traction / proof
s = slide()
header(s, "Proof", "Born from a platform that already ran nationally")
bullets(s, [
    ("MNAIR — Make Nigeria AI Ready:", "our national AI-literacy platform and LearnStack's first flagship curriculum."),
    ("The engine is not a prototype:", "multi-tenant platform built, end-to-end tested (authoring → enrollment → quiz → verified certificate), MNAIR imports with one command."),
    ("Institutional demand signal:", "NITDA–NYSC digital literacy ambassador program actively seeking platform partners (2026)."),
], size=16)
stat_row(s, [("175", "lessons authored"), ("5", "Nigerian languages"), ("254", "pilot learners"), ("11", "states reached")], y=Inches(4.75))
notes(s, "Pilot: 254 WhatsApp users across 11 states; update completion-rate figure from analytics before presenting.")

# 6 — Why now
s = slide()
header(s, "Why now", "The window is open — and the incumbents aren't looking here")
bullets(s, [
    ("Government push.", "Nigeria's National AI Strategy is in implementation; NITDA is recruiting platform partners for nationwide digital literacy."),
    ("WhatsApp-first learning is exploding,", "but no incumbent LMS delivers structured curricula with media over WhatsApp, SMS, or USSD."),
    ("Global e-learning spend keeps compounding", "(LMS market est. $20B+, high-teens CAGR) while African training businesses still run on PDFs and group chats."),
    ("The differentiation window is real:", "channel delivery + white-labeling + content portability is a structural combination, not a feature toggle."),
])

# 7 — Competition
s = slide()
header(s, "Competition", "The combination is the moat")
table_slide(s, ["", "LearnStack", "Teachable/Thinkific", "Moodle", "TalentLMS"], [
    ["True white-label", "✓ core design", "partial, premium", "✓ with effort", "partial"],
    ["Run many academies (operator model)", "✓", "✗", "✗", "limited"],
    ["WhatsApp / SMS / USSD delivery", "✓ shipped", "✗", "✗", "✗"],
    ["Curriculum portability (full JSON)", "✓", "✗", "partial", "✗"],
    ["Verifiable certificates", "✓", "add-on", "plugin", "✓"],
    ["Ops burden on customer", "none", "none", "high", "none"],
], col_widths=[Inches(3.9), Inches(2.2), Inches(2.4), Inches(1.7), Inches(1.7)])

# 8 — Business model
s = slide()
header(s, "Business model", "Per-tenant SaaS, priced on learners — never on content")
table_slide(s, ["Tier", "Target", "Price", "Includes"], [
    ["Starter", "Solo trainers", "$29/mo", "1 live course, 100 learners, subdomain"],
    ["Growth", "Academies", "$99/mo", "Unlimited courses, 1,000 learners, custom domain"],
    ["Institution", "Gov / NGO / corporate", "$499+/mo", "Unlimited learners, SLA, onboarding, channels"],
    ["Sovereign", "National programs", "Custom", "Isolated deployment, operator training"],
], col_widths=[Inches(1.8), Inches(2.8), Inches(1.6), Inches(5.7)])
tf = box(s, Inches(0.7), Inches(4.7), Inches(11.9), Inches(2.2))
para(tf, "Marginal cost per tenant is near zero (shared infrastructure). High-margin add-ons: "
         "curriculum migration services (half-automated already), certificate volume for verifiers, "
         "and later a revenue share on paid-course checkout. Paystack-first for NGN reality; Stripe for the rest.",
     15, DARK, first=True)

# 9 — Go-to-market
s = slide()
header(s, "Go-to-market", "Land where the pain is sharpest, expand where the budgets are")
bullets(s, [
    ("1 · Training academies & bootcamps (now).", "West African trainers running on WhatsApp groups and PDFs — direct sales, 5 design partners first."),
    ("2 · Government & NGO programs.", "The NITDA–NYSC ambassador model generalizes: every agency becomes a tenant. MNAIR is the working reference."),
    ("3 · Corporate L&D and compliance.", "Certification with a public audit trail; SCORM import unlocks migration deals."),
    ("4 · Edtech founders (self-serve).", "The Teachable segment, minus the brand tax and lock-in."),
])

# 10 — Roadmap
s = slide()
header(s, "Roadmap", "Shipped platform → first revenue → compounding moat")
bullets(s, [
    ("Shipped:", "multi-tenant white-labeling · block authoring with video · WhatsApp/SMS/USSD delivery with simulator · quizzes · verifiable certificates · JSON import/export · Docker deploy · test suite."),
    ("Next 90 days:", "Paystack + Stripe billing · S3/CDN media · per-tenant analytics · custom-domain automation · hosted drip cadences."),
    ("Two quarters:", "cohorts & scheduled programs · multi-language variants · WhatsApp interactive buttons · AI-assisted lesson drafting."),
    ("Four quarters:", "SCORM/xAPI import · voice/IVR channel · curriculum marketplace with revenue share."),
], size=15.5)

# 11 — The ask
s = slide(NAVY)
rect(s, 0, 0, W, Inches(0.18), GREEN)
tf = box(s, Inches(0.9), Inches(0.8), Inches(11.5), Inches(1.2))
para(tf, "THE ASK", 13, GREEN, bold=True, first=True)
para(tf, "Pre-seed to turn a shipped platform into a revenue machine", 30, WHITE, bold=True)
tf = box(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(3.6))
para(tf, "Raising: $250K pre-seed  (12–15 months runway)", 20, WHITE, bold=True, first=True, space_after=14)
para(tf, "•  40% — engineering: billing, S3/CDN media, analytics, cohort scheduling", 16, RGBColor(0xD9, 0xE4, 0xDE), space_after=8)
para(tf, "•  30% — go-to-market: 5 design partners → 50 paying tenants; NITDA partnership", 16, RGBColor(0xD9, 0xE4, 0xDE), space_after=8)
para(tf, "•  20% — infrastructure & provider costs (WhatsApp Business, Africa's Talking)", 16, RGBColor(0xD9, 0xE4, 0xDE), space_after=8)
para(tf, "•  10% — operations & compliance (NDPA)", 16, RGBColor(0xD9, 0xE4, 0xDE), space_after=14)
para(tf, "Milestone at raise+12mo: 50 tenants · 25,000 learners · $15K MRR · 2 institutional contracts",
     17, WHITE, bold=True)
notes(s, "EDIT: amount, allocation, and milestones are founder placeholders — set your real numbers before presenting.")

# 12 — Close
s = slide()
rect(s, 0, Inches(6.9), W, Inches(0.6), NAVY)
tf = box(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(3))
para(tf, "Everyone has something to teach.", 34, NAVY, bold=True, first=True, space_after=6)
para(tf, "We give them the academy.", 34, GREEN, bold=True, space_after=24)
para(tf, "LearnStack — white-label curriculum delivery on web, WhatsApp, SMS, and USSD.", 16, GREY)
notes(s, "Add contact details / demo link (http://localhost:3003/?tenant=demo for live demos).")

out = Path(__file__).resolve().parent.parent / "PITCH_DECK.pptx"
prs.save(out)
print(out)
